"""Slide deck rendering service for generating professional presentations."""

from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Optional
from uuid import UUID

from ...core.config import settings
from ...models.rendering import (
    ComparisonColumn,
    ContentType,
    DeckRenderRequest,
    ExportFormat,
    MetricItem,
    PricingTier,
    RenderResult,
    RenderStatus,
    Slide,
    SlideContent,
    SlideLayout,
    TeamMember,
    TimelineItem,
    WebDeckConfig,
    WebDeckRenderRequest,
    WebDeckResult,
)
from .brand_styler import BrandStyler


class DeckRenderer:
    """Renders professional slide decks from content JSON."""

    def __init__(self, brand_config: Optional[dict] = None):
        """Initialize the deck renderer."""
        self.styler = BrandStyler(brand_config) if brand_config else BrandStyler()
        self._ensure_output_dir()

    def _ensure_output_dir(self) -> None:
        """Ensure output directory exists."""
        settings.output_dir.mkdir(parents=True, exist_ok=True)

    def _get_layout_class(self, layout: SlideLayout) -> str:
        """Get CSS class for slide layout."""
        return f"layout-{layout.value.replace('_', '-')}"

    def _render_slide_html(self, slide: Slide, index: int, total: int, show_numbers: bool) -> str:
        """Render a single slide to HTML."""
        content = slide.content
        layout_class = self._get_layout_class(content.layout)

        # Build slide HTML
        parts = [
            f"<div class='slide {layout_class}' id='slide-{slide.id}' data-transition='{slide.transition}'>",
        ]

        # Logo
        logo_path = self.styler.get_logo_path()
        if logo_path and content.layout != SlideLayout.TITLE:
            parts.append(f"<img src='file://{logo_path}' alt='Logo' class='slide-logo'>")

        # Background
        if content.background_image:
            parts.append(f"<div class='slide-bg' style='background-image: url({content.background_image})'></div>")

        # Render based on layout
        if content.layout == SlideLayout.TITLE:
            parts.extend(self._render_title_slide(content))
        elif content.layout == SlideLayout.QUOTE:
            parts.extend(self._render_quote_slide(content))
        elif content.layout == SlideLayout.METRICS:
            parts.extend(self._render_metrics_slide(content))
        elif content.layout == SlideLayout.TEAM:
            parts.extend(self._render_team_slide(content))
        elif content.layout == SlideLayout.PRICING:
            parts.extend(self._render_pricing_slide(content))
        elif content.layout == SlideLayout.TIMELINE:
            parts.extend(self._render_timeline_slide(content))
        elif content.layout == SlideLayout.COMPARISON:
            parts.extend(self._render_comparison_slide(content))
        elif content.layout == SlideLayout.CTA:
            parts.extend(self._render_cta_slide(content))
        elif content.layout in (SlideLayout.IMAGE_LEFT, SlideLayout.IMAGE_RIGHT):
            parts.extend(self._render_image_text_slide(content))
        elif content.layout == SlideLayout.FULL_IMAGE:
            parts.extend(self._render_full_image_slide(content))
        elif content.layout == SlideLayout.TWO_COLUMN:
            parts.extend(self._render_two_column_slide(content))
        else:
            parts.extend(self._render_standard_slide(content))

        # Slide number
        if show_numbers:
            parts.append(f"<div class='slide-number'>{index + 1} / {total}</div>")

        parts.append("</div>")

        return "\n".join(parts)

    def _render_title_slide(self, content: SlideContent) -> list[str]:
        """Render title slide layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")
        if content.subtitle:
            parts.append(f"<p class='slide-subtitle'>{content.subtitle}</p>")
        return parts

    def _render_standard_slide(self, content: SlideContent) -> list[str]:
        """Render standard title + content slide."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")
        if content.subtitle:
            parts.append(f"<p class='slide-subtitle'>{content.subtitle}</p>")

        parts.append("<div class='slide-body'>")

        if content.body:
            for block in content.body:
                parts.append(f"<p>{block.content}</p>")

        if content.bullets:
            list_tag = "ol" if content.bullets.style == "numbered" else "ul"
            parts.append(f"<{list_tag}>")
            for item in content.bullets.items:
                parts.append(f"<li>{item}</li>")
            parts.append(f"</{list_tag}>")

        if content.image:
            parts.append(
                f"<img src='{content.image.url}' alt='{content.image.alt_text}' class='slide-image'>"
            )

        parts.append("</div>")
        return parts

    def _render_quote_slide(self, content: SlideContent) -> list[str]:
        """Render quote layout."""
        parts = []
        if content.quote:
            parts.append(f"<blockquote class='quote-text'>{content.quote}</blockquote>")
        if content.quote_author:
            parts.append(f"<p class='quote-author'>— {content.quote_author}</p>")
        return parts

    def _render_metrics_slide(self, content: SlideContent) -> list[str]:
        """Render metrics layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        if content.metrics:
            parts.append("<div class='metrics-grid'>")
            for metric in content.metrics:
                trend_class = f" {metric.trend}" if metric.trend else ""
                parts.append(
                    f"<div class='metric-card'>"
                    f"<div class='metric-value'>{metric.value}</div>"
                    f"<div class='metric-label'>{metric.label}</div>"
                )
                if metric.trend:
                    arrow = "↑" if metric.trend == "up" else "↓" if metric.trend == "down" else "→"
                    parts.append(f"<div class='metric-trend{trend_class}'>{arrow}</div>")
                if metric.description:
                    parts.append(f"<div class='metric-desc'>{metric.description}</div>")
                parts.append("</div>")
            parts.append("</div>")
        return parts

    def _render_team_slide(self, content: SlideContent) -> list[str]:
        """Render team layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        if content.team:
            parts.append("<div class='team-grid'>")
            for member in content.team:
                parts.append("<div class='team-member'>")
                if member.image_url:
                    parts.append(f"<img src='{member.image_url}' alt='{member.name}'>")
                else:
                    parts.append("<div class='member-avatar'></div>")
                parts.append(f"<div class='member-name'>{member.name}</div>")
                parts.append(f"<div class='member-role'>{member.role}</div>")
                if member.bio:
                    parts.append(f"<div class='member-bio'>{member.bio}</div>")
                parts.append("</div>")
            parts.append("</div>")
        return parts

    def _render_pricing_slide(self, content: SlideContent) -> list[str]:
        """Render pricing layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        if content.pricing:
            parts.append("<div class='pricing-grid'>")
            for tier in content.pricing:
                highlight_class = " highlighted" if tier.highlighted else ""
                parts.append(f"<div class='pricing-tier{highlight_class}'>")
                parts.append(f"<div class='tier-name'>{tier.name}</div>")
                parts.append(f"<div class='tier-price'>{tier.price}</div>")
                if tier.description:
                    parts.append(f"<div class='tier-desc'>{tier.description}</div>")
                if tier.features:
                    parts.append("<ul class='tier-features'>")
                    for feature in tier.features:
                        parts.append(f"<li>{feature}</li>")
                    parts.append("</ul>")
                parts.append(f"<button class='tier-cta'>{tier.cta_text}</button>")
                parts.append("</div>")
            parts.append("</div>")
        return parts

    def _render_timeline_slide(self, content: SlideContent) -> list[str]:
        """Render timeline layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        if content.timeline:
            parts.append("<div class='timeline'>")
            for item in content.timeline:
                parts.append("<div class='timeline-item'>")
                parts.append(f"<div class='timeline-date'>{item.date}</div>")
                parts.append("<div class='timeline-content'>")
                parts.append(f"<div class='timeline-title'>{item.title}</div>")
                if item.description:
                    parts.append(f"<div class='timeline-description'>{item.description}</div>")
                parts.append("</div>")
                parts.append("</div>")
            parts.append("</div>")
        return parts

    def _render_comparison_slide(self, content: SlideContent) -> list[str]:
        """Render comparison table layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        if content.comparison and content.comparison_rows:
            parts.append("<table class='comparison-table'>")
            parts.append("<thead><tr><th></th>")
            for col in content.comparison:
                highlight = " class='highlighted'" if col.highlighted else ""
                parts.append(f"<th{highlight}>{col.header}</th>")
            parts.append("</tr></thead>")

            parts.append("<tbody>")
            for i, row_label in enumerate(content.comparison_rows):
                parts.append(f"<tr><td>{row_label}</td>")
                for col in content.comparison:
                    if i < len(col.values):
                        highlight = " class='highlighted'" if col.highlighted else ""
                        parts.append(f"<td{highlight}>{col.values[i]}</td>")
                parts.append("</tr>")
            parts.append("</tbody></table>")
        return parts

    def _render_cta_slide(self, content: SlideContent) -> list[str]:
        """Render call-to-action layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")
        if content.cta_text:
            url = content.cta_url or "#"
            parts.append(f"<a href='{url}' class='cta-button'>{content.cta_text}</a>")
        return parts

    def _render_image_text_slide(self, content: SlideContent) -> list[str]:
        """Render image + text side by side."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        parts.append("<div class='slide-content'>")

        # Image
        if content.image:
            parts.append(
                f"<div class='slide-image-container'>"
                f"<img src='{content.image.url}' alt='{content.image.alt_text}' class='slide-image'>"
                f"</div>"
            )

        # Text content
        parts.append("<div class='slide-text'>")
        if content.body:
            for block in content.body:
                parts.append(f"<p>{block.content}</p>")
        if content.bullets:
            list_tag = "ol" if content.bullets.style == "numbered" else "ul"
            parts.append(f"<{list_tag}>")
            for item in content.bullets.items:
                parts.append(f"<li>{item}</li>")
            parts.append(f"</{list_tag}>")
        parts.append("</div>")

        parts.append("</div>")
        return parts

    def _render_full_image_slide(self, content: SlideContent) -> list[str]:
        """Render full-bleed image slide."""
        parts = []
        if content.image:
            parts.append(f"<img src='{content.image.url}' alt='{content.image.alt_text}' class='slide-image'>")
            parts.append("<div class='slide-overlay'></div>")
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")
        return parts

    def _render_two_column_slide(self, content: SlideContent) -> list[str]:
        """Render two-column layout."""
        parts = []
        if content.title:
            parts.append(f"<h1 class='slide-title'>{content.title}</h1>")

        parts.append("<div class='slide-content'>")

        # Split body content between columns if available
        if content.body and len(content.body) >= 2:
            parts.append("<div class='column'>")
            parts.append(f"<p>{content.body[0].content}</p>")
            parts.append("</div>")
            parts.append("<div class='column'>")
            for block in content.body[1:]:
                parts.append(f"<p>{block.content}</p>")
            parts.append("</div>")
        elif content.body:
            parts.append("<div class='column'>")
            for block in content.body:
                parts.append(f"<p>{block.content}</p>")
            parts.append("</div>")

        parts.append("</div>")
        return parts

    def generate_html(self, request: DeckRenderRequest) -> str:
        """Generate HTML for the entire deck."""
        slides_html = []
        total = len(request.slides)

        for i, slide in enumerate(request.slides):
            slides_html.append(
                self._render_slide_html(slide, i, total, request.show_slide_numbers)
            )

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{request.title}</title>
    <style>
    {self.styler.get_slide_styles()}
    </style>
</head>
<body>
    <div class="deck">
        {"".join(slides_html)}
    </div>
</body>
</html>"""
        return html

    async def render_to_html(self, request: DeckRenderRequest) -> RenderResult:
        """Render deck as HTML file."""
        try:
            html_content = self.generate_html(request)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c if c.isalnum() else "_" for c in request.title)[:50]
            filename = f"{safe_title}_{timestamp}.html"
            output_path = settings.output_dir / filename

            output_path.write_text(html_content, encoding="utf-8")

            return RenderResult(
                id=request.id,
                status=RenderStatus.COMPLETED,
                content_type=request.content_type,
                format=ExportFormat.HTML,
                completed_at=datetime.utcnow(),
                file_path=str(output_path),
                file_url=f"/downloads/{filename}",
                file_size=output_path.stat().st_size,
                page_count=len(request.slides),
                metadata={"title": request.title},
            )
        except Exception as e:
            return RenderResult(
                id=request.id,
                status=RenderStatus.FAILED,
                content_type=request.content_type,
                format=ExportFormat.HTML,
                error_message=str(e),
            )

    async def render_to_pptx(self, request: DeckRenderRequest) -> RenderResult:
        """Render deck as PowerPoint file."""
        try:
            from pptx import Presentation
            from pptx.dml.color import RgbColor
            from pptx.util import Inches, Pt

            # Create presentation
            prs = Presentation()

            # Set slide dimensions (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for slide_data in request.slides:
                content = slide_data.content
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # Add title if present
                if content.title:
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
                    )
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = content.title
                    title_para.font.size = Pt(44)
                    title_para.font.bold = True
                    # Convert hex to RGB
                    hex_color = self.styler.config.primary_color.lstrip("#")
                    title_para.font.color.rgb = RgbColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16),
                    )

                # Add subtitle if present
                if content.subtitle:
                    sub_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5)
                    )
                    sub_frame = sub_box.text_frame
                    sub_para = sub_frame.paragraphs[0]
                    sub_para.text = content.subtitle
                    sub_para.font.size = Pt(24)
                    hex_color = self.styler.config.secondary_color.lstrip("#")
                    sub_para.font.color.rgb = RgbColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16),
                    )

                # Add body content
                body_top = 2.2 if content.subtitle else 1.8

                if content.body:
                    body_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(body_top), Inches(12.333), Inches(4.5)
                    )
                    body_frame = body_box.text_frame
                    body_frame.word_wrap = True

                    for i, block in enumerate(content.body):
                        if i == 0:
                            para = body_frame.paragraphs[0]
                        else:
                            para = body_frame.add_paragraph()
                        para.text = block.content
                        para.font.size = Pt(20)
                        para.space_after = Pt(12)

                # Add bullets
                if content.bullets:
                    bullet_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(body_top + (1 if content.body else 0)),
                        Inches(12.333), Inches(4)
                    )
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.word_wrap = True

                    for i, item in enumerate(content.bullets.items):
                        if i == 0:
                            para = bullet_frame.paragraphs[0]
                        else:
                            para = bullet_frame.add_paragraph()
                        para.text = item
                        para.font.size = Pt(18)
                        para.level = 0
                        para.space_after = Pt(8)

                # Add quote for quote slides
                if content.quote:
                    quote_box = slide.shapes.add_textbox(
                        Inches(1), Inches(2.5), Inches(11.333), Inches(2)
                    )
                    quote_frame = quote_box.text_frame
                    quote_para = quote_frame.paragraphs[0]
                    quote_para.text = f'"{content.quote}"'
                    quote_para.font.size = Pt(28)
                    quote_para.font.italic = True

                    if content.quote_author:
                        author_para = quote_frame.add_paragraph()
                        author_para.text = f"— {content.quote_author}"
                        author_para.font.size = Pt(20)

            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c if c.isalnum() else "_" for c in request.title)[:50]
            filename = f"{safe_title}_{timestamp}.pptx"
            output_path = settings.output_dir / filename

            prs.save(str(output_path))

            return RenderResult(
                id=request.id,
                status=RenderStatus.COMPLETED,
                content_type=request.content_type,
                format=ExportFormat.PPTX,
                completed_at=datetime.utcnow(),
                file_path=str(output_path),
                file_url=f"/downloads/{filename}",
                file_size=output_path.stat().st_size,
                page_count=len(request.slides),
                metadata={"title": request.title},
            )

        except Exception as e:
            return RenderResult(
                id=request.id,
                status=RenderStatus.FAILED,
                content_type=request.content_type,
                format=ExportFormat.PPTX,
                error_message=str(e),
            )

    async def render(self, request: DeckRenderRequest) -> RenderResult:
        """Render deck to the requested format."""
        if request.format == ExportFormat.HTML:
            return await self.render_to_html(request)
        elif request.format == ExportFormat.PPTX:
            return await self.render_to_pptx(request)
        elif request.format == ExportFormat.PDF:
            # Render HTML then convert to PDF
            html_content = self.generate_html(request)
            return await self._html_to_pdf(request, html_content)
        else:
            return RenderResult(
                id=request.id,
                status=RenderStatus.FAILED,
                content_type=request.content_type,
                format=request.format,
                error_message=f"Unsupported format: {request.format}",
            )

    async def _html_to_pdf(self, request: DeckRenderRequest, html_content: str) -> RenderResult:
        """Convert HTML deck to PDF."""
        try:
            from weasyprint import HTML, CSS

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = "".join(c if c.isalnum() else "_" for c in request.title)[:50]
            filename = f"{safe_title}_{timestamp}.pdf"
            output_path = settings.output_dir / filename

            # Add print styles for slide-like pages
            print_css = CSS(string=f"""
                @page {{
                    size: {settings.deck_width}px {settings.deck_height}px;
                    margin: 0;
                }}
                .slide {{
                    page-break-after: always;
                    width: {settings.deck_width}px;
                    height: {settings.deck_height}px;
                }}
            """)

            html = HTML(string=html_content, base_url=str(settings.base_dir))
            html.write_pdf(output_path, stylesheets=[print_css])

            return RenderResult(
                id=request.id,
                status=RenderStatus.COMPLETED,
                content_type=request.content_type,
                format=ExportFormat.PDF,
                completed_at=datetime.utcnow(),
                file_path=str(output_path),
                file_url=f"/downloads/{filename}",
                file_size=output_path.stat().st_size,
                page_count=len(request.slides),
                metadata={"title": request.title},
            )
        except Exception as e:
            return RenderResult(
                id=request.id,
                status=RenderStatus.FAILED,
                content_type=request.content_type,
                format=ExportFormat.PDF,
                error_message=str(e),
            )
